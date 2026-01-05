import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_slide_images(pptx_path, slide_number, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation(pptx_path)
    
    # Slide numbers are 1-indexed, but list is 0-indexed
    slide = prs.slides[slide_number - 1]
    
    print(f"Extracting images from slide {slide_number}...")
    
    image_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.content_type.split('/')[1]
            filename = f"Priya_slide10_{image_count}.{ext}"
            file_path = os.path.join(output_folder, filename)
            
            with open(file_path, 'wb') as f:
                f.write(image.blob)
            
            print(f"[OK] Extracted: {filename} (size: {len(image.blob)} bytes)")
            image_count += 1

file_path = r"C:\Users\arieldre\Downloads\Persona based comm June 2025 (3).pptx"
output_folder = r"C:\Users\arieldre\.cursor\agentic project\persona-pulse\public\images\personas"
extract_slide_images(file_path, 10, output_folder)

