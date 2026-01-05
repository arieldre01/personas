import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_persona_images(pptx_path, output_folder="persona_images"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation(pptx_path)
    
    # List of expected personas to look for in text
    target_personas = ["Maya", "Ben", "Oliver", "Priya", "Anna", "Sahil", "Ido", "Alex"]
    
    print(f"Scanning {len(prs.slides)} slides for personas...")

    for i, slide in enumerate(prs.slides):
        # 1. Get all text on the slide to identify the persona
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        
        # 2. Match text to a persona name
        found_name = None
        for name in target_personas:
            if name in slide_text:
                found_name = name
                break
        
        # 3. If a persona is found, extract the pictures on that slide
        if found_name:
            image_count = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    # Get the image file extension (e.g., png, jpg)
                    ext = image.content_type.split('/')[1]
                    
                    # Create a filename (e.g., Maya.jpeg)
                    # If multiple images exist (like icons), we number them: Maya_0.jpeg
                    filename = f"{found_name}_{image_count}.{ext}"
                    file_path = os.path.join(output_folder, filename)
                    
                    with open(file_path, 'wb') as f:
                        f.write(image.blob)
                    
                    print(f"[OK] Extracted: {filename}")
                    image_count += 1

# Update this path to where your file is located
file_path = r"C:\Users\arieldre\Downloads\Persona based comm June 2025 (3).pptx"
output_folder = r"C:\Users\arieldre\.cursor\agentic project\persona-pulse\public\images\personas"
extract_persona_images(file_path, output_folder)

